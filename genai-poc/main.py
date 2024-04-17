import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain_community.llms import OpenAI
from langchain_openai import OpenAI
import pandas as pd
from io import BytesIO
from pptx import Presentation
import os
from docx import Document


class SessionState:
    def __init__(self, **kwargs):
        self.__dict__.update(kwargs)


def process_documents(documents_folder):
    documents = []
    for file_name in os.listdir(documents_folder):
        file_path = os.path.join(documents_folder, file_name)
        if file_path.endswith(".pdf"):
            pdf_reader = PdfReader(file_path)
            for page in pdf_reader.pages:
                documents.append(page.extract_text())
        elif file_path.endswith(".txt"):
            with open(file_path, "r") as file:
                documents.append(file.read())
        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                documents.append(paragraph.text)
        elif file_path.endswith(".xlsx"):
            xls = pd.ExcelFile(file_path)
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                for col in df.columns:
                    documents.extend(df[col].dropna().astype(str).tolist())
        elif file_path.endswith(".pptx"):
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        documents.append(shape.text)

    text = "\n".join(documents)
    if not text:
        return None

    char_text_splitter = CharacterTextSplitter(separator="\n", chunk_size=1000, chunk_overlap=200, length_function=len)
    text_chunks = char_text_splitter.split_text(text)
    if not text_chunks:
        return None

    embeddings = OpenAIEmbeddings(openai_api_key="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")
    docsearch = FAISS.from_texts(text_chunks, embeddings)
    return docsearch


def answer_query(query, docsearch, chain):
    docs = docsearch.similarity_search(query)
    response = chain.run(input_documents=docs, question=query)
    # response = chain.invoke(input_documents=docs, question=query)
    return response


def main():
    st.header("Innovative Bot 🤖")
    session_state = SessionState(documents_folder=None, docsearch=None, response="", is_admin=False)
    # Authentication mechanism
    user_type = st.radio("Select user type:", ["Admin", "User"])
    chain = None  # Initialize chain variable to None
    if user_type == "Admin":
        admin_email = st.text_input("Enter admin email:")
        admin_password = st.text_input("Enter admin password:", type="password")
        login_button = st.button("Login")

        if admin_password == "Manju@300602" and admin_email == "manju123@gmail.com":
            session_state.is_admin = True
            st.success("Login successfully")

        else:
            st.error("Incorrect password and email. Please try again.")
            return

        if session_state.is_admin:
            # Admin functionality
            uploaded_files = st.file_uploader("Upload your documents", accept_multiple_files=True,
                                              type=["pdf", "txt", "docx", "xlsx", "pptx"])
            documents_folder = "upload_documents"
            if not os.path.exists(documents_folder):
                os.makedirs(documents_folder)
            for file in uploaded_files:
                with open(os.path.join(documents_folder, file.name), "wb") as f:
                    f.write(file.getvalue())

            session_state.documents_folder = documents_folder
            session_state.docsearch = process_documents(documents_folder)
            chain = load_qa_chain(OpenAI(api_key="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"),
                                  chain_type="stuff")

            # Display uploaded documents
            st.subheader("Uploaded Documents")
            selected_documents = []
            for file_name in os.listdir(documents_folder):
                selected = st.checkbox(file_name)
                if selected:
                    selected_documents.append(file_name)

            # Delete selected documents
            if st.button("Delete Selected Documents"):
                for file_name in selected_documents:
                    file_path = os.path.join(documents_folder, file_name)
                    os.remove(file_path)
                st.write("Documents deleted successfully.")

    else:
        # Normal user functionality
        documents_folder = "upload_documents"  # Here we are using the folder where admin will upload the documents
        if not os.path.exists(documents_folder):
            st.write("No documents uploaded yet. Please wait for the admin to upload documents.")
            return

        session_state.documents_folder = documents_folder
        session_state.docsearch = process_documents(documents_folder)

        # Display uploaded documents
        st.subheader("Uploaded Documents")
        for file_name in os.listdir(session_state.documents_folder):
            st.write(file_name)

    # Ensure chain is loaded if documents are uploaded
    if session_state.docsearch:
        chain = load_qa_chain(OpenAI(api_key="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"), chain_type="stuff")

    # AI assistant chat interface
    st.sidebar.header("AI Assistant 🚀")
    st.sidebar.subheader("Conversation History")
    chat_history = st.sidebar.empty()
    user_input = st.sidebar.text_input("Query:", key="user_input")
    if st.sidebar.button("Send"):
        if chain:  # Check if chain is not None
            response = answer_query(user_input, session_state.docsearch, chain)
            st.sidebar.text("Chat Bot:")
            st.sidebar.write(response)
        else:
            st.sidebar.warning("AI Assistant is not available. Please wait for the admin to upload documents.")


if __name__ == '__main__':
    main()



